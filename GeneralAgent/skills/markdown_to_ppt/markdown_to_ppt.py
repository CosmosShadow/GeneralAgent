def markdown_to_ppt(markdown_file_path):
    """
    convert markdown file to ppt. return ppt file path.
    """
    import re
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # 读取Markdown文件
    with open(markdown_file_path, 'r', encoding='utf-8') as file:
        markdown_content = file.read()

    # 创建PPT文档
    prs = Presentation()

    # 定义一些样式参数
    title_font_size = Pt(44)
    subtitle_font_size = Pt(32)
    content_font_size = Pt(24)
    bullet_font_size = Pt(18)
    max_content_lines = 12  # 假设每页最多显示12行内容

    # 解析Markdown内容并提取一级标题作为PPT标题页
    title_match = re.search(r'^# (.*)', markdown_content, re.MULTILINE)
    if title_match:
        # 添加标题页
        slide_layout = prs.slide_layouts[0]  # 选择标题页布局
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        # 设置标题页内容
        title.text = title_match.group(1).strip()
        subtitle.text = "Subtitle or Additional Information"

        # 从Markdown内容中移除一级标题
        markdown_content = re.sub(r'^# .*\n\n', '', markdown_content, count=1, flags=re.MULTILINE)

    # 解析Markdown内容
    sections = re.split(r'\n(?=## )', markdown_content)  # 按二级标题分割内容
    for section in sections:
        # 提取标题
        title_match = re.search(r'^(#+)\s*(.*)', section, re.MULTILINE)
        if not title_match:
            continue

        # 创建新的幻灯片
        slide_layout = prs.slide_layouts[1]  # 选择带标题的内容页布局
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]

        # 设置标题
        title_shape.text = title_match.group(2).strip()
        title_shape.text_frame.paragraphs[0].font.size = title_font_size

        # 添加内容
        content_lines = section.split('\n')[1:]  # 去掉标题行
        content_paragraphs = [line.strip() for line in content_lines if line.strip() and not line.strip().startswith('#')]

        for i, paragraph in enumerate(content_paragraphs):
            if i > 0 and i % max_content_lines == 0:
                # 如果内容太多，分割到新的幻灯片
                slide = prs.slides.add_slide(slide_layout)
                content_shape = slide.placeholders[1]

            # 处理列表项
            if paragraph.startswith('- '):
                paragraph = paragraph[2:].strip()  # 移除列表项符号
                level = 1
            else:
                level = 0

            # 添加段落
            p = content_shape.text_frame.add_paragraph()
            p.text = paragraph
            p.level = level
            p.font.size = bullet_font_size if level > 0 else content_font_size
            p.alignment = PP_ALIGN.LEFT
            
    # 保存PPT文档
    from GeneralAgent import skills
    ppt_file_path = skills.unique_name() + '.pptx'
    prs.save(ppt_file_path)

    return ppt_file_path


def test_markdown_to_ppt():
    import os
    markdown_file_path = os.path.join(os.path.dirname(__file__), 'bid_plan.md')
    ppt_file_path = markdown_to_ppt(markdown_file_path)

if __name__ == '__main__':
    test_markdown_to_ppt()