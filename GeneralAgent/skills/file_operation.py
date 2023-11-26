    
def read_pdf_pages(file_path):
    """Read the pdf file and return a list of strings on each page of the pdf"""
    """读取pdf文件，返回pdf每页字符串的列表"""
    import fitz
    doc = fitz.open(file_path)
    documents = []
    for page in doc:
        documents.append(page.get_text())
    return documents
    
def read_word_pages(file_path):
    """Read the word file and return a list of word paragraph strings"""
    """读取word文件，返回word段落字符串的列表"""
    # https://zhuanlan.zhihu.com/p/146363527
    from docx import Document
    # 打开文档
    document = Document(file_path)
    # 读取标题、段落、列表内容
    ps = [ paragraph.text for paragraph in document.paragraphs]
    return ps

def read_ppt(file_path):
    import pptx
    prs = pptx.Presentation(file_path)
    documents = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                documents.append(shape.text)
    return '\n'.join(documents)


def read_file_content(file_path):
    """return content of txt, md, pdf, docx file"""
    # 支持file_path的类型包括txt、md、pdf、docx
    if file_path.endswith('.pdf'):
        return ' '.join(read_pdf_pages(file_path))
    elif file_path.endswith('.docx'):
        return ' '.join(read_word_pages(file_path))
    elif file_path.endswith('.ppt') or file_path.endswith('.pptx'):
        return read_ppt(file_path)
    else:
        # 默认当做文本文件
        with open(file_path, 'r', encoding='utf-8') as f:
            return '\n'.join(f.readlines())

def write_file_content(file_path, content):
    """write content to txt, md"""
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(content)
