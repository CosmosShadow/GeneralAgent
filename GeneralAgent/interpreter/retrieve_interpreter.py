# read the document and can retrieve the information
import re
from .interpreter import Interpreter
import chromadb
import logging
from GeneralAgent.llm import embedding_batch, num_tokens_from_string

def _split_text(text, max_len):
    """
    Splits a given text into segments of up to max_len characters, using punctuation marks as delimiters.
    """
    import re
    segments = []
    current_segment = ""
    for sentence in re.split(r"[，。,.;；]", text):
        # 本身就是大于max_len的句子
        if len(sentence) > max_len:
            for i in range(0, len(sentence), max_len):
                segments.append(sentence[i:i+max_len])
        else:
            if len(current_segment + sentence) <= max_len:
                current_segment += sentence
            else:
                segments.append(current_segment)
                current_segment = sentence
    if current_segment:
        segments.append(current_segment)
    return segments

def _merge_short_strings(paragraphs, min_length=100):
    merged_paragraphs = []
    current_string = ""
    for p in paragraphs:
        current_string += ' ' + p
        if len(current_string) < min_length:
            continue
        else:
            merged_paragraphs.append(current_string)
            current_string = ''
    if len(current_string) > 0:
        merged_paragraphs.append(current_string)
    return merged_paragraphs

# 文本转换为段落
def _text_to_paragraphs(text, max_len=510):
    if len(text) == 0:
        return []
    
    # 粗略分段
    paragraph_list = text.split("\n")

    # 合并短句
    paragraph_list = _merge_short_strings(paragraph_list)

    results = []

    # 再精细分段
    for paragraph in paragraph_list:
        if len(paragraph) == 0:
            continue
        if len(paragraph) * 2 <= max_len:
            results.append(paragraph)
        else:
            results += _split_text(paragraph, max_len)
    results = [x.strip() for x in results if len(x.strip()) > 0]
    return results

def read_pdf(file_path):
    import fitz
    doc = fitz.open(file_path)
    documents = []
    for page in doc:
        documents.append(page.get_text())
    return '\n'.join(documents)

def read_ppt(file_path):
    import pptx
    prs = pptx.Presentation(file_path)
    documents = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                documents.append(shape.text)
    return '\n'.join(documents)

def read_docx(file_path):
    import docx
    doc = docx.Document(file_path)
    documents = []
    for para in doc.paragraphs:
        documents.append(para.text)
    return '\n'.join(documents)


class RetrieveInterpreter(Interpreter):
    def __init__(self, serialize_path='./read_data/', prompt_max_length=1000, useful_msg_count=2) -> None:
        self.prompt_max_length = prompt_max_length
        self.useful_msg_count = useful_msg_count
        self.client = chromadb.PersistentClient(path=serialize_path)
        self.collection = self.client.get_or_create_collection(name="read", metadata={"hnsw:space": "cosine"})

    def prompt(self, messages) -> str:
        # when collection is empty, return empty string
        if self.collection.count() == 0:
            return ''
        
        querys = []
        for x in messages[-self.useful_msg_count:]:
            querys += _text_to_paragraphs(x['content'])
        query_embeddings = embedding_batch(querys)
        result = self.collection.query(
            query_embeddings=query_embeddings,
            n_results=2,
        )
        # extract distances and documents
        distances = [x for z in result['distances'] for x in z]
        documents = [x for z in result['documents'] for x in z]

        # sort distances and documents by distance
        sorted_docs = sorted(list(zip(distances, documents)), key=lambda x: x[0])

        # filter documents with distance < 100
        documents = [x for d, x in sorted_docs if d < 100]
        # texts token count should less than prompt_max_length
        texts = []
        texts_token_count = 0
        for x in documents:
            if texts_token_count + num_tokens_from_string(x) > self.prompt_max_length:
                break
            texts.append(x)
        return '\n'.join(texts)

    @property
    def match_template(self):
        return '```read\n(.*?)\n```'
    
    def parse(self, string):
        information = []
        pattern = re.compile(self.match_template, re.DOTALL)
        match = pattern.search(string)
        assert match is not None
        file_paths = match.group(1).strip().split('\n')
        for file_path in file_paths:
            paragraphs = []
            if file_path.endswith('.pdf'):
                paragraphs = _text_to_paragraphs(read_pdf(file_path))
            elif file_path.endswith('.pptx'):
                paragraphs = _text_to_paragraphs(read_ppt(file_path))
            elif file_path.endswith('.docx'):
                paragraphs = _text_to_paragraphs(read_docx(file_path))
            else:
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                    paragraphs = _text_to_paragraphs(text)
                except Exception as e:
                    logging.exception(e)
                    information.append(f'read the content of file {file_path} fails: ' + str(e))
            if len(paragraphs) > 0:
                information.append(f'The content of file {file_path} is: ' + '\n'.join(paragraphs)[:100] + '\n......')
            embeddings = embedding_batch(paragraphs)
            logging.debug(paragraphs[:2])
            logging.debug(embeddings[:2])
            self.collection.add(
                documents=paragraphs,
                embeddings=embeddings,
                metadatas=[{'file_path': file_path} for _ in paragraphs],
                ids=[file_path+str(i) for i in range(len(paragraphs))],
            )
        stop = False
        if string.replace(match.group(0), '').strip() == '':
            stop = True
        info = '\n'.join(information)
        string = f'\n{string}\n```{info}```\n'
        return string, stop